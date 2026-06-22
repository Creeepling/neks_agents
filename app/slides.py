import io
import os
from typing import List, Optional, Any, Dict
from pptx import Presentation as PPTXPresentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pydantic import BaseModel
import instructor
from google import genai

from app.config import settings

# ---------------------------------------------------------------------------
# LLM Schema for Presentation
# ---------------------------------------------------------------------------
class PresentationSlide(BaseModel):
    title: str
    bullets: List[str]

class PresentationData(BaseModel):
    city_analysis: Optional[PresentationSlide] = None
    location_analysis: Optional[PresentationSlide] = None
    competitors: Optional[PresentationSlide] = None
    generated_concepts: Optional[PresentationSlide] = None
    finalized_concepts: List[PresentationSlide] = []


def extract_presentation_data(prop_data: Dict[str, Any]) -> PresentationData:
    """Uses Gemini to summarize the raw property data into structured slides."""
    client = instructor.from_genai(
        client=genai.Client(api_key=settings.GEMINI_API_KEY),
        use_async=False
    )
    
    prompt = f"""
    You are an expert commercial real estate analyst.
    Your task is to convert the following raw JSON data into a clean, professional slide deck structure.
    Keep bullet points concise (max 10-15 words each).
    
    Rules for extraction:
    - If there is not enough data for a slide, leave it null.
    - 'competitors' should list the direct competitors.
    - 'generated_concepts' should list the initial concepts brainstormed.
    - 'finalized_concepts' MUST have one slide for EVERY finalized concept in the data. Make sure to include the proposed tenants in the bullets!
    
    Property Data:
    {prop_data}
    """
    
    return client.chat.completions.create(
        model="gemini-2.5-flash",
        response_model=PresentationData,
        messages=[{"role": "user", "content": prompt}]
    )

# ---------------------------------------------------------------------------
# PPTX Generation
# ---------------------------------------------------------------------------
def _apply_theme(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 15, 16) # #0F0F10

def _format_title(title_shape):
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(250, 67, 31) # #FA431F
    title_shape.text_frame.paragraphs[0].font.bold = True

def _format_bullets(body_shape):
    for paragraph in body_shape.text_frame.paragraphs:
        paragraph.font.color.rgb = RGBColor(255, 255, 255) # White
        paragraph.font.size = Pt(20)

def generate_pptx(prop_name: str, prop_address: str, data: PresentationData) -> io.BytesIO:
    prs = PPTXPresentation()
    
    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    _apply_theme(slide)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = prop_name
    _format_title(title)
    
    subtitle.text = prop_address or "Коммерческая недвижимость"
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(147, 147, 147) # #939393
    
    # Add Logo if exists
    logo_path = os.path.join(os.path.dirname(__file__), "assets", "logo.png")
    if os.path.exists(logo_path):
        prs.slides[0].shapes.add_picture(logo_path, Inches(0.5), Inches(0.5), height=Inches(0.5))

    # Helper for adding standard slides
    def add_bullet_slide(slide_data: PresentationSlide):
        if not slide_data or not slide_data.bullets:
            return
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        _apply_theme(slide)
        
        if os.path.exists(logo_path):
            slide.shapes.add_picture(logo_path, Inches(0.5), Inches(0.3), height=Inches(0.3))
            
        title_shape = slide.shapes.title
        title_shape.text = slide_data.title
        _format_title(title_shape)
        
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.text = slide_data.bullets[0]
        
        for bullet in slide_data.bullets[1:]:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0
            
        _format_bullets(body_shape)

    if data.city_analysis:
        add_bullet_slide(data.city_analysis)
    if data.location_analysis:
        add_bullet_slide(data.location_analysis)
    if data.competitors:
        add_bullet_slide(data.competitors)
    if data.generated_concepts:
        add_bullet_slide(data.generated_concepts)
    
    for concept in data.finalized_concepts:
        add_bullet_slide(concept)
        
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output
